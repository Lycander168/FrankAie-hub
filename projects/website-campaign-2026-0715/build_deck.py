#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
官網全站活動簡報產生器 — 2026 暑假 ｜ Wokyis × LYCANDER
依 README.md 內容產出品牌化 campaign-deck.pptx。可重跑、可改。
用法: python3 build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- 品牌色 ----------
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)   # 科技深藍（主）
ORANGE = RGBColor(0xFF, 0x7A, 0x29)   # 暑假活力橙（accent）
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF8)   # 淺灰底
GREY   = RGBColor(0x5A, 0x63, 0x72)   # 次要文字
DARK   = RGBColor(0x22, 0x27, 0x30)   # 內文深色
TEAL   = RGBColor(0x18, 0x9A, 0xB4)   # 輔助

EMU_W, EMU_H = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

FONT = "Microsoft JhengHei"   # 微軟正黑（缺則由 PowerPoint 替代）


# ---------- 共用小工具 ----------
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) or (text, size, bold, color, space_after)"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        text, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        sa = ln[4] if len(ln) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
    return tb

def chip(s, x, y, w, text, fill, txtcolor=WHITE, size=12):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = True; f.color.rgb = txtcolor; f.name = FONT
    return sp

def header(s, kicker, title):
    rect(s, 0, 0, EMU_W, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), EMU_W, Inches(0.08), ORANGE)
    textbox(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.4),
            [(kicker, 12, True, ORANGE)])
    textbox(s, Inches(0.6), Inches(0.44), Inches(12), Inches(0.7),
            [(title, 26, True, WHITE)])

def footer(s, idx):
    textbox(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.4),
            [("官網全站活動 ｜ Wokyis × LYCANDER ｜ 2026/7/15", 9, False, GREY)])
    textbox(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.4),
            [(str(idx), 9, True, GREY)], align=PP_ALIGN.RIGHT)

def placeholder(s, x, y, w, h, label="商品圖\n（置入官方圖）"):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = LIGHT
    sp.line.color.rgb = GREY; sp.line.width = Pt(1)
    sp.line.dash_style = 2 if hasattr(sp.line, "dash_style") else None
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    for i, t in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t
        f = r.font; f.size = Pt(11); f.color.rgb = GREY; f.name = FONT
    return sp

def table(s, x, y, w, h, data, col_w=None, header_fill=NAVY, fs=11, hfs=11,
          zebra=True, align_first_left=True):
    rows, cols = len(data), len(data[0])
    gtbl = s.shapes.add_table(rows, cols, x, y, w, h).table
    gtbl.first_row = True; gtbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for j, cw in enumerate(col_w):
            gtbl.columns[j].width = Emu(int(w * cw / total))
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = gtbl.cell(i, j)
            c.margin_left = Pt(5); c.margin_right = Pt(5)
            c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (j == 0 and align_first_left) else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            f = r.font; f.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
                f.size = Pt(hfs); f.bold = True; f.color.rgb = WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if (not zebra or i % 2 == 1) else LIGHT
                f.size = Pt(fs); f.bold = False; f.color.rgb = DARK
    return gtbl


# ============================================================
# 1. 封面
# ============================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, NAVY)
rect(s, 0, Inches(5.0), EMU_W, Inches(0.14), ORANGE)
# 裝飾圓
c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.4), Inches(-1.2), Inches(4.2), Inches(4.2))
c1.fill.solid(); c1.fill.fore_color.rgb = TEAL; c1.line.fill.background(); c1.shadow.inherit = False
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.6), Inches(4.6), Inches(3.4), Inches(3.4))
c2.fill.solid(); c2.fill.fore_color.rgb = ORANGE; c2.line.fill.background(); c2.shadow.inherit = False
chip(s, Inches(0.9), Inches(1.5), Inches(2.6), "會員招募中 ｜ 2026 暑假", ORANGE)
textbox(s, Inches(0.9), Inches(2.2), Inches(11), Inches(2.2),
        [("官網全站活動企劃", 48, True, WHITE, 6),
         ("Wokyis 關聯配件　×　LYCANDER 暑假活動", 24, True, RGBColor(0xCF,0xDA,0xEC), 6)])
textbox(s, Inches(0.9), Inches(5.4), Inches(11), Inches(1.4),
        [("上線日期：2026 / 7 / 15　｜　暑假主檔 7/15–8/31", 16, True, WHITE, 4),
         ("交付對象：行銷公司　｜　炫輝國際 Shinetech / LYCANDER", 13, False, RGBColor(0xB7,0xC4,0xDA), 4)])

# ============================================================
# 2. 活動總覽
# ============================================================
s = slide(); header(s, "OVERVIEW", "活動總覽 — 一頁看懂")
cards = [
    ("主題", "會員招募中", "暑假入會衝刺", ORANGE),
    ("檔期", "7/15 開跑", "主檔 7/15–8/31", TEAL),
    ("主軸 A", "Wokyis 關聯配件", "M5 擴充座新品首發", NAVY),
    ("主軸 B", "LYCANDER 暑假活動", "行動辦公 / 開學季", NAVY),
]
cx = Inches(0.6)
for k, v, d, col in cards:
    rect(s, cx, Inches(1.5), Inches(2.95), Inches(1.9), LIGHT)
    rect(s, cx, Inches(1.5), Inches(2.95), Inches(0.12), col)
    textbox(s, cx+Inches(0.2), Inches(1.75), Inches(2.6), Inches(1.6),
            [(k, 12, True, col, 4), (v, 17, True, DARK, 4), (d, 11, False, GREY, 2)])
    cx += Inches(3.13)
# 三機制
textbox(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.4), [("三大機制", 15, True, NAVY)])
mech = [
    ("① 會員招募", "入會即享簡單入會禮 ＋ 會員專屬活動價"),
    ("② 滿額贈", "單筆滿 NT$1,000 送 YOYOISLES AirTag 保護套（市值 NT$190）"),
    ("③ 活動組合", "桌面生產力組 / 行動辦公組 / 開學季組 / 配件加價購"),
]
mx = Inches(0.6)
for t, d in mech:
    rect(s, mx, Inches(4.2), Inches(3.95), Inches(1.5), NAVY)
    textbox(s, mx+Inches(0.22), Inches(4.4), Inches(3.6), Inches(1.2),
            [(t, 15, True, ORANGE, 6), (d, 12, False, WHITE, 2)])
    mx += Inches(4.13)
textbox(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.7),
        [("目標 KPI：新入會數 ｜ 客單價 AOV ｜ 轉化率 ｜ 滿額贈核銷率 ｜ 組合連帶率", 13, True, TEAL)])
footer(s, 2)

# ============================================================
# 3. 檔期節奏
# ============================================================
s = slide(); header(s, "TIMELINE", "檔期節奏")
steps = [
    ("預熱", "7/8–7/14", "會員招募預告\nEDM/社群倒數\n名單蒐集", TEAL),
    ("開跑", "7/15", "全站上線\n首頁 KV 換版\n滿額贈啟動", ORANGE),
    ("主檔", "7/15–8/31", "Wokyis 新品首發\nLYCANDER 暑假組合\n會員價持續", NAVY),
    ("收尾", "8/25–8/31", "最後一週催單\n開學季加碼\n贈品庫存倒數", GREY),
]
# 時間軸線
rect(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.06), ORANGE)
cx = Inches(0.8)
for name, period, detail, col in steps:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, Inches(1.85), Inches(0.36), Inches(0.36))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.color.rgb = WHITE; dot.line.width = Pt(2)
    dot.shadow.inherit = False
    rect(s, cx, Inches(2.5), Inches(2.7), Inches(2.5), LIGHT)
    rect(s, cx, Inches(2.5), Inches(2.7), Inches(0.5), col)
    textbox(s, cx, Inches(2.55), Inches(2.7), Inches(0.45),
            [(f"{name}　{period}", 13, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, cx+Inches(0.2), Inches(3.15), Inches(2.4), Inches(1.7),
            [(d, 12, False, DARK, 4) for d in detail.split("\n")])
    cx += Inches(2.95)
footer(s, 3)

# ============================================================
# 4. 機制 1：會員招募
# ============================================================
s = slide(); header(s, "MECHANIC 1", "會員招募中 — 簡單入會禮")
textbox(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
        [("入會即享（二擇一，建議由業主拍板）", 15, True, NAVY)])
# 兩方案卡
rect(s, Inches(0.6), Inches(1.95), Inches(5.9), Inches(1.95), LIGHT)
rect(s, Inches(0.6), Inches(1.95), Inches(0.14), Inches(1.95), ORANGE)
textbox(s, Inches(0.9), Inches(2.15), Inches(5.4), Inches(1.6),
        [("方案 A ｜ NT$100 購物金", 17, True, ORANGE, 6),
         ("首購折抵，門檻低、拉新最快。", 13, False, DARK, 4),
         ("適合衝「新入會數」KPI。", 12, False, GREY, 2)])
rect(s, Inches(6.8), Inches(1.95), Inches(5.9), Inches(1.95), LIGHT)
rect(s, Inches(6.8), Inches(1.95), Inches(0.14), Inches(1.95), TEAL)
textbox(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(1.6),
        [("方案 B ｜ YOYOISLES AirTag 保護套", 17, True, TEAL, 6),
         ("實物有感（市值 NT$190），與品牌調性合。", 13, False, DARK, 4),
         ("適合提升入會「體感價值」。", 12, False, GREY, 2)])
# 會員專屬 + CTA
textbox(s, Inches(0.6), Inches(4.15), Inches(12), Inches(0.4),
        [("會員專屬 & 招募導流", 15, True, NAVY)])
table(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(1.7),
      [["項目", "內容"],
       ["會員專屬", "活動期間全站會員價、組合包會員加碼"],
       ["招募 CTA", "首頁 banner ｜ 商品頁入會浮層 ｜ 結帳前「登入多省 NT$○」"],
       ["導流管道", "EDM（7/8 起預熱）｜ 社群限動倒數 ｜ 既有客回購名單"]],
      col_w=[2, 8], fs=12, hfs=12)
footer(s, 4)

# ============================================================
# 5. 機制 2：滿額贈
# ============================================================
s = slide(); header(s, "MECHANIC 2", "滿額贈 — 單筆滿千送 AirTag 保護套")
placeholder(s, Inches(9.4), Inches(1.6), Inches(3.3), Inches(3.3), "YOYOISLES\nAir Stars\nAirTag 保護套圖")
chip(s, Inches(0.6), Inches(1.6), Inches(3.6), "市值 NT$190 ｜ YOYOISLES Air Stars", ORANGE)
textbox(s, Inches(0.6), Inches(2.2), Inches(8.5), Inches(1.0),
        [("單筆結帳滿 NT$1,000", 30, True, NAVY, 4),
         ("即送 星型 EVA AirTag 保護套 ×1（4 色）", 18, True, ORANGE, 4)])
table(s, Inches(0.6), Inches(3.5), Inches(8.5), Inches(2.5),
      [["規則", "說明"],
       ["門檻", "單筆結帳金額滿 NT$1,000（含）"],
       ["贈品", "YOYOISLES Air Stars AirTag 保護套 ×1，市值 NT$190"],
       ["數量上限", "每筆 1 個、送完為止（庫存上限【需業主確認】）"],
       ["疊加", "可與會員價／組合價疊加；贈品不折現"],
       ["加購誘導", "門檻貼近單品客單 → 文案主打「再加一條線材即達標」"]],
      col_w=[2, 8], fs=11.5, hfs=12)
footer(s, 5)

# ============================================================
# 6. 主軸 A：Wokyis 關聯配件
# ============================================================
s = slide(); header(s, "TRACK A ｜ 新品首發", "Wokyis 關聯配件 — Mac mini 桌面生態核心")
placeholder(s, Inches(0.6), Inches(1.55), Inches(3.6), Inches(3.4), "Wokyis M5\n擴充座圖")
textbox(s, Inches(4.45), Inches(1.55), Inches(8.2), Inches(1.6),
        [("Wokyis M5 迷你螢幕擴充座（hero 新品）", 18, True, NAVY, 6),
         ("復古外型 × 現代效能：5吋 IPS 副螢幕＋M.2 NVMe SSD 擴充＋10Gbps 多孔，", 13, False, DARK, 2),
         ("把 Mac mini 變成桌面主角。10Gbps / 80Gbps TB5 兩版本。", 13, False, DARK, 2)])
textbox(s, Inches(4.45), Inches(3.25), Inches(8.2), Inches(0.4),
        [("為何「關聯配件」打法最有效", 14, True, ORANGE)])
for i, (t, d) in enumerate([
        ("螢幕外接", "Wokyis M5 ＋ LYCANDER HDMI 2.1 8K 線 → 一線到位"),
        ("桌面充電", "＋ OLIKA W3 三合一磁吸 → 手機/手錶/耳機歸位"),
        ("輸入升級", "＋ iPad/鍵盤周邊 → 完整桌面生產力"),
        ("拉高客單", "單買 vs 組合 → AOV 提升、滿額贈門檻自然達標")]):
    y = Inches(3.7 + i*0.72)
    rect(s, Inches(4.45), y, Inches(8.2), Inches(0.62), LIGHT)
    rect(s, Inches(4.45), y, Inches(0.1), Inches(0.62), TEAL)
    textbox(s, Inches(4.7), y+Inches(0.06), Inches(7.9), Inches(0.5),
            [(f"{t}　—　{d}", 12, False, DARK)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# ============================================================
# 7. Wokyis 組合
# ============================================================
s = slide(); header(s, "TRACK A ｜ 組合", "Wokyis 關聯配件組合")
table(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.0),
      [["組合", "內容", "原價(推估)", "組合價", "省"],
       ["桌面生產力組", "M5 10G ＋ HDMI 8K 線 ＋ OLIKA W3", "NT$8,370", "NT$7,490", "NT$880"],
       ["創作者旗艦組", "M5 80G TB5 ＋ Mac mini 9合1 底座 ＋ 100W 線", "NT$13,980", "NT$12,490", "NT$1,490"]],
      col_w=[2.3, 5.2, 1.7, 1.6, 1.2], fs=12, hfs=12)
rect(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(1.2), LIGHT)
rect(s, Inches(0.6), Inches(3.9), Inches(0.14), Inches(1.2), ORANGE)
textbox(s, Inches(0.9), Inches(4.05), Inches(11.6), Inches(1.0),
        [("組合策略", 14, True, ORANGE, 4),
         ("以 Wokyis M5 為錨點帶動 LYCANDER 周邊連帶銷售；兩組合皆 >NT$1,000，自動觸發滿額贈，", 12.5, False, DARK, 2),
         ("讓客人「買整套更划算又有禮」，提升客單價與新品曝光。", 12.5, False, DARK, 2)])
textbox(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
        [("＊售價為市場推估，正式以官網後台 SKU／採購報價覆核。", 11, False, GREY)])
footer(s, 7)

# ============================================================
# 8. 主軸 B：LYCANDER 暑假活動
# ============================================================
s = slide(); header(s, "TRACK B ｜ 暑假主題", "LYCANDER 暑假活動 — 行動辦公 / 開學季")
cats = [
    ("筆電包", "Contrastin 支架包\nStiksels 拼接包", "行動辦公"),
    ("iPad 鍵盤", "Combo Lite\nHALFTER", "開學季生產力"),
    ("桌面充電", "OLIKA W3\n100W 編織線", "桌面整潔"),
    ("影音線材", "HDMI 2.1 8K\n極細認證線", "外接必備"),
]
cx = Inches(0.6)
for t, items, tag in cats:
    rect(s, cx, Inches(1.6), Inches(2.95), Inches(2.7), LIGHT)
    rect(s, cx, Inches(1.6), Inches(2.95), Inches(0.55), NAVY)
    textbox(s, cx, Inches(1.66), Inches(2.95), Inches(0.45),
            [(t, 15, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    placeholder(s, cx+Inches(0.55), Inches(2.3), Inches(1.85), Inches(1.0), "主圖")
    textbox(s, cx+Inches(0.2), Inches(3.4), Inches(2.6), Inches(0.9),
            [(ln, 11.5, False, DARK, 2) for ln in items.split("\n")])
    chip(s, cx+Inches(0.55), Inches(3.95), Inches(1.85), tag, ORANGE, size=10)
    cx += Inches(3.13)
textbox(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.4),
        [("暑假情境包裝", 14, True, ORANGE, 4),
         ("以「行動辦公／開學季」故事線包裝既有核心品類，搭配限時組合與加價購；", 12.5, False, DARK, 2),
         ("主打族群：學生、文書/SOHO、Mac 使用者。", 12.5, False, DARK, 2)])
footer(s, 8)

# ============================================================
# 9. LYCANDER 組合
# ============================================================
s = slide(); header(s, "TRACK B ｜ 組合", "LYCANDER 暑假活動組合")
table(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.3),
      [["組合", "內容", "原價(推估)", "組合價", "省"],
       ["行動辦公組", "Contrastin 筆電包 ＋ HDMI 8K 線 ＋ OLIKA W3", "NT$4,160", "NT$3,690", "NT$470"],
       ["開學季 iPad 組", "Combo Lite 鍵盤 ＋ 180° 收納包（＋AirTag 套贈）", "NT$3,270", "NT$2,790", "NT$480＋贈"],
       ["配件加價購", "任一主商品 +100W 線 / +180° 收納包", "—", "+NT$390 起", "—"]],
      col_w=[2.2, 5.3, 1.7, 1.6, 1.3], fs=12, hfs=12)
rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(1.0), LIGHT)
rect(s, Inches(0.6), Inches(4.2), Inches(0.14), Inches(1.0), TEAL)
textbox(s, Inches(0.9), Inches(4.33), Inches(11.6), Inches(0.8),
        [("加價購是滿額贈的最佳推手", 13.5, True, TEAL, 4),
         ("結帳前推「+NT$390 線材」既湊滿千得贈品，又提升件數與連帶率。", 12.5, False, DARK, 2)])
textbox(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
        [("＊售價為市場推估，正式以官網後台 SKU／採購報價覆核。", 11, False, GREY)])
footer(s, 9)

# ============================================================
# 10. 建議上架品項總表
# ============================================================
s = slide(); header(s, "ASSORTMENT", "建議上架品項總表")
data = [["品項", "品牌", "角色", "建議售價", "活動價", "關聯/說明"],
    ["M5 迷你螢幕擴充座 10Gbps", "Wokyis", "🆕新品", "5,490*", "4,990", "Mac mini 桌面核心"],
    ["M5 擴充座 80Gbps TB5", "Wokyis", "🆕新品", "10,900*", "9,990", "高階/創作者"],
    ["HDMI 2.1 8K 認證線 1.5M", "LYCANDER", "⭐主推", "690", "590", "外接螢幕必備/好湊單"],
    ["OLIKA W3 三合一磁吸充電", "LYCANDER", "⭐主推", "2,190", "1,890", "桌面充電主力"],
    ["Combo Lite iPad 鍵盤殼", "LYCANDER", "⭐主推", "2,680", "2,380", "開學季主打"],
    ["Contrastin 磁吸筆電包", "LYCANDER", "⭐主推", "1,280*", "1,080", "行動辦公"],
    ["Mac mini 9合1 擴充底座", "LYCANDER", "⭐主推", "2,490*", "2,190", "與 Wokyis 互補"],
    ["USB-C 100W 編織線【新建議】", "LYCANDER", "🔁衝量", "590*", "490", "耗材回購/湊單"],
    ["M.2 NVMe SSD 外接盒【新建議】", "（新）", "🔁衝量", "690*", "590", "Mac 桌面儲存"],
    ["螢幕增高架/理線【新建議】", "（新）", "🔁衝量", "890*", "790", "桌面美學/提客單"],
    ["Air Stars AirTag 保護套", "YOYOISLES", "🎁贈品", "市值190", "贈品", "滿額贈/入會禮"]]
table(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.0),
      data, col_w=[3.3, 1.5, 1.2, 1.3, 1.1, 3.0], fs=10.5, hfs=10.5)
textbox(s, Inches(0.4), Inches(6.55), Inches(12), Inches(0.4),
        [("＊標 * 或【新建議】之售價為推估，待覆核；單位 NT$。", 10.5, False, GREY)])
footer(s, 10)

# ============================================================
# 11. 活動組合總表
# ============================================================
s = slide(); header(s, "BUNDLES", "活動組合總表")
table(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(3.3),
      [["組合", "內容", "原價*", "組合價", "省"],
       ["Wokyis 桌面生產力組", "M5 10G ＋ HDMI 8K ＋ OLIKA W3", "8,370", "7,490", "880"],
       ["Wokyis 創作者旗艦組", "M5 80G TB5 ＋ Mac mini 底座 ＋ 100W 線", "13,980", "12,490", "1,490"],
       ["LYCANDER 行動辦公組", "Contrastin 包 ＋ HDMI 8K ＋ OLIKA W3", "4,160", "3,690", "470"],
       ["開學季 iPad 組", "Combo Lite 鍵盤 ＋ 180° 收納包（＋贈AirTag）", "3,270", "2,790", "480＋贈"],
       ["配件加價購", "主商品 +100W 線 / +180° 收納包", "—", "+390 起", "—"]],
      col_w=[2.6, 5.4, 1.3, 1.4, 1.3], fs=11.5, hfs=11.5)
rect(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.0), NAVY)
textbox(s, Inches(0.8), Inches(5.25), Inches(11.8), Inches(0.8),
        [("全組合皆 ≥ NT$1,000 → 自動觸發滿額贈，疊加感強化「買整套更划算又有禮」。", 13, True, WHITE, 4),
         ("＊單位 NT$，售價為推估，正式以後台覆核。", 10.5, False, RGBColor(0xB7,0xC4,0xDA), 2)])
footer(s, 11)

# ============================================================
# 12. 視覺與文案方向
# ============================================================
s = slide(); header(s, "CREATIVE", "視覺與文案方向")
textbox(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4), [("主視覺 Slogan（候選）", 14, True, NAVY)])
for i, t in enumerate(["「暑假升級你的桌面 — 會員招募中」",
                       "「Mac mini 一桌到位 × 暑假行動辦公」"]):
    rect(s, Inches(0.6+i*6.1), Inches(1.95), Inches(5.85), Inches(0.95), LIGHT)
    rect(s, Inches(0.6+i*6.1), Inches(1.95), Inches(0.12), Inches(0.95), ORANGE)
    textbox(s, Inches(0.85+i*6.1), Inches(2.1), Inches(5.4), Inches(0.7),
            [(t, 15, True, DARK)], anchor=MSO_ANCHOR.MIDDLE)
# 主色
textbox(s, Inches(0.6), Inches(3.15), Inches(12), Inches(0.4), [("品牌主色", 14, True, NAVY)])
for i, (name, col, hexv) in enumerate([("科技深藍", NAVY, "#1B2A4A"), ("活力橙", ORANGE, "#FF7A29"),
                                       ("輔助青", TEAL, "#189AB4"), ("淺灰底", LIGHT, "#F2F4F8")]):
    x = Inches(0.6+i*3.05)
    sw = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.6), Inches(2.85), Inches(0.7))
    sw.fill.solid(); sw.fill.fore_color.rgb = col; sw.line.color.rgb = GREY; sw.line.width = Pt(0.5)
    sw.shadow.inherit = False
    tc = WHITE if name != "淺灰底" else DARK
    textbox(s, x+Inches(0.15), Inches(3.7), Inches(2.6), Inches(0.5),
            [(f"{name}  {hexv}", 12, True, tc)], anchor=MSO_ANCHOR.MIDDLE)
# Banner 文案
textbox(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4), [("Banner 文案", 14, True, NAVY)])
table(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.4),
      [["位置", "文案"],
       ["首頁 KV", "7/15 全站開跑｜入會即享好禮｜滿千送 AirTag 保護套"],
       ["商品頁浮層", "登入會員再省，還送暑假好禮"],
       ["素材清單", "首頁 KV／會員招募 banner／滿額贈 banner／Wokyis 新品主圖／組合主圖／社群倒數"]],
      col_w=[2, 9], fs=12, hfs=12)
footer(s, 12)

# ============================================================
# 13. 執行 checklist
# ============================================================
s = slide(); header(s, "EXECUTION", "官網執行 Checklist")
checks = [
    ("1 商品上架", "依品項表建立/更新 SKU、價格、活動價；新建議品標待覆核"),
    ("2 機制設定", "入會禮（擇一）、滿額贈（滿千自動加贈、庫存上限）、組合與加價購"),
    ("3 視覺換版", "首頁 KV、兩支活動 banner、商品頁浮層，7/15 00:00 排程"),
    ("4 會員導流", "EDM 預熱（7/8 起）、社群倒數、結帳前會員提示"),
    ("5 Listing", "把完整文案貼入對應商品（標題/五點/描述/SEO）"),
    ("6 QA", "滿額贈觸發、會員價、組合價、行動版顯示、贈品庫存扣減"),
    ("7 雙通路一致", "官網確認後同步蝦皮 shinetech88888"),
]
y = Inches(1.5)
for t, d in checks:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(0.66), LIGHT)
    rect(s, Inches(0.6), y, Inches(2.3), Inches(0.66), NAVY)
    textbox(s, Inches(0.7), y, Inches(2.2), Inches(0.66),
            [(t, 12.5, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(3.1), y, Inches(9.4), Inches(0.66),
            [(d, 12, False, DARK)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.74)
footer(s, 13)

# ============================================================
# 14. KPI
# ============================================================
s = slide(); header(s, "KPI", "KPI 與成效追蹤")
kpis = [("新入會數", "招募成效", ORANGE), ("客單價 AOV", "組合/加購帶動", TEAL),
        ("轉化率", "活動頁→結帳", NAVY), ("滿額贈核銷率", "門檻有效性", ORANGE),
        ("組合連帶率", "套組吸引力", TEAL), ("Wokyis 新品銷量", "首發動能", NAVY)]
cx, cy = Inches(0.6), Inches(1.6)
for i, (t, d, col) in enumerate(kpis):
    x = Inches(0.6 + (i % 3)*4.13)
    yy = Inches(1.6 + (i // 3)*2.2)
    rect(s, x, yy, Inches(3.95), Inches(1.9), LIGHT)
    rect(s, x, yy, Inches(3.95), Inches(0.12), col)
    textbox(s, x+Inches(0.25), yy+Inches(0.35), Inches(3.5), Inches(1.4),
            [(t, 19, True, DARK, 6), (d, 12, False, GREY, 2)])
textbox(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.6),
        [("建議：每週檢視，主檔中段（8/1）做一次成效複盤與素材/組合微調。", 13, True, TEAL)])
footer(s, 14)

# ============================================================
# 15–19. 附錄 A：完整 Listing 文案
# ============================================================
listings = [
    ("Wokyis M5 迷你螢幕擴充座（hero 新品）",
     "Wokyis M5 Mac mini 迷你螢幕擴充座｜5吋IPS螢幕+SSD擴充+10Gbps USB Hub 復古桌面神器",
     ["Mac mini 一秒變身：5吋 IPS 副螢幕，顯示系統資訊/時鐘/監控，復古造型滿足質感控。",
      "SSD 擴充：內建 M.2 NVMe 槽（10Gbps），桌面儲存一次補滿。",
      "一座多孔：4×USB-A ＋ USB-C 10Gbps ＋ SD/microSD ＋ 3.5mm 耳機孔。",
      "廣相容：Mac mini M4/M2/M1、MacBook Air/Pro、Windows 筆電皆可用。",
      "一桌到位：搭 LYCANDER HDMI 8K 線與 OLIKA W3 充電，打造整潔高效桌面。"],
     "Mac mini 擴充座,Wokyis,M5,迷你螢幕,NVMe SSD 擴充,USB Hub,Mac mini M4,副螢幕"),
    ("LYCANDER HDMI 2.1 8K 極細認證線 1.5M",
     "LYCANDER HDMI 2.1 8K 極細認證影音傳輸線 1.5M｜48Gbps 8K@60Hz 認證線 Mac/PS5/螢幕",
     ["8K@60Hz / 4K@120Hz：HDMI 2.1 滿頻 48Gbps，畫面清晰流暢不掉格。",
      "極細好收：細線徑桌面不雜亂，搭 Wokyis M5/Mac mini 外接螢幕剛好。",
      "官方認證：相容 Mac、PS5、Switch、螢幕、投影。",
      "耐折耐用：強化接頭與線身，插拔更安心。",
      "滿額好搭：NT$690 一條，湊滿千送 AirTag 套無痛達標。"],
     "HDMI 2.1,8K 傳輸線,48Gbps,認證線,Mac 外接螢幕,PS5 HDMI,1.5M,LYCANDER"),
    ("LYCANDER OLIKA W3 三合一磁吸無線充電器",
     "LYCANDER OLIKA W3 三合一磁吸無線充電器｜iPhone+Apple Watch+AirPods 桌面快充座",
     ["三合一一次充：iPhone 磁吸 ＋ Apple Watch ＋ AirPods 一座搞定。",
      "磁吸對位：MagSafe 相容，放上即充不對位煩惱。",
      "桌面整潔：取代多條線材，與 Wokyis/Mac 桌面美學一致。",
      "快充穩定：智慧分配、過充保護。",
      "行動辦公必備：桌面、出差、床頭一座通用。"],
     "三合一無線充電,MagSafe,磁吸充電座,Apple Watch 充電,AirPods 充電,OLIKA W3"),
    ("LYCANDER Combo Lite 觸控鍵盤 iPad 保護殼",
     "LYCANDER Combo Lite iPad 觸控鍵盤磁吸可拆保護殼｜巧控觸控板 注音鍵盤 開學季生產力",
     ["鍵盤＋觸控板：iPad 秒變輕筆電，報告打字更快。",
      "磁吸可拆：看劇純平板、打字接鍵盤。",
      "注音在地化：印刷注音鍵位，直覺好上手。",
      "多角度支撐：站立辦公、追劇皆穩固。",
      "開學季首選：學生、文書族暑假升級裝備。"],
     "iPad 鍵盤,巧控鍵盤,觸控板鍵盤,iPad 保護殼,注音鍵盤,Combo Lite,LYCANDER"),
    ("YOYOISLES Air Stars AirTag 保護套（滿額贈／入會禮）",
     "YOYOISLES Air Stars 星型 AirTag 保護套｜EVA 輕量防護 門卡磁扣 安全無毒 4色",
     ["星型設計：辨識度高、療癒造型，AirTag 一秒變吊飾。",
      "EVA 輕量防護：精密模具包覆，防刮防撞不阻擋訊號。",
      "掛繩孔好攜帶：鑰匙、背包、寵物項圈隨掛隨找。",
      "安全無毒：材質安心，送禮自用皆宜。",
      "4 色可選：黑／紫／黃／綠。"],
     "AirTag 保護套,AirTag 殼,星型,EVA,鑰匙圈,YOYOISLES,Air Stars,定位器套"),
]
for idx, (name, title, bullets, seo) in enumerate(listings):
    s = slide()
    header(s, f"APPENDIX A ｜ Listing {idx+1}/5", name)
    placeholder(s, Inches(10.0), Inches(1.5), Inches(2.8), Inches(2.6), "商品圖")
    chip(s, Inches(0.6), Inches(1.5), Inches(1.4), "商品標題", NAVY)
    rect(s, Inches(0.6), Inches(2.0), Inches(9.1), Inches(0.95), LIGHT)
    textbox(s, Inches(0.8), Inches(2.1), Inches(8.7), Inches(0.8),
            [(title, 13, True, DARK)], anchor=MSO_ANCHOR.MIDDLE)
    chip(s, Inches(0.6), Inches(3.15), Inches(1.4), "賣點五點", ORANGE)
    bl = []
    for i, b in enumerate(bullets):
        bl.append((f"{i+1}. {b}", 12.5, False, DARK, 7))
    textbox(s, Inches(0.6), Inches(3.65), Inches(12.1), Inches(2.4), bl)
    chip(s, Inches(0.6), Inches(6.05), Inches(1.8), "SEO 後台關鍵字", TEAL)
    textbox(s, Inches(2.6), Inches(6.08), Inches(10), Inches(0.6),
            [(seo, 11, False, GREY)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 15+idx)

# ============================================================
# 20. 附錄 B：來源與前提
# ============================================================
s = slide(); header(s, "APPENDIX B", "來源與前提")
textbox(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4), [("待覆核（正式上線前確認）", 14, True, ORANGE)])
textbox(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(1.5),
        [("• 所有台幣售價／活動價／毛利、贈品庫存上限 → 以官網後台 SKU 與採購報價覆核。", 12, False, DARK, 5),
         ("• Wokyis、YOYOISLES 是否確定上架官網／蝦皮 shinetech88888（目前未上架）。", 12, False, DARK, 5),
         ("• 會員入會禮最終方案（NT$100 購物金 vs YOYOISLES AirTag 套）。", 12, False, DARK, 5),
         ("• 商品圖：簡報佔位框由行銷公司置入官方商品圖。", 12, False, DARK, 5)])
textbox(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4), [("公開來源", 14, True, NAVY)])
textbox(s, Inches(0.6), Inches(4.05), Inches(12.1), Inches(2.6),
        [("• Wokyis：wokyis.com（M5 10Gbps / 80Gbps）、Amazon B0FPB67QBV、Kickstarter「Wokyis M5」。", 11.5, False, DARK, 5),
         ("• YOYOISLES：yoyoisles.com、Amazon B0F3NLL3R5、SunSky EDA007793101A。", 11.5, False, DARK, 5),
         ("• LYCANDER：lycander.tw、lycander.tw/hdmi、Yahoo 購物（HDMI 8K、OLIKA W3）、PChome（Combo Lite）、STUDIO A、momo。", 11.5, False, DARK, 5),
         ("• 蝦皮賣場：shopee.tw/shinetech88888（炫輝國際，133 件商品）。", 11.5, False, DARK, 5),
         ("• 方法論：FrankAie-hub — marketing-team / listing-optimizer / pricing-calculator / consumers-taiwan / seo-keyword-expander。", 11.5, False, GREY, 5)])
footer(s, 20)

# ---------- 輸出 ----------
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "campaign-deck.pptx")
prs.save(out)
print(f"OK: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
