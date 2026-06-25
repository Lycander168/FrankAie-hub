# -*- coding: utf-8 -*-
"""
三檔齊發團購提報簡報生成器（8/1 on 檔）
產出：三檔齊發團購提報_0801.pptx（16:9，繁體中文，約 12 頁）

對象：團購業者 ama shop ── 提供 Wokyis / LYCANDER / Sharge 三檔商品與活動資訊。

設計原則
- 商業條件（建議售價／團購價／拆帳%／MOQ／出貨日）一律留「待填 TBD」欄位，不臆測數字。
- 產品圖／KV 以虛線佔位框呈現，由業務後置（嘖嘖頁面無法程式抓圖）。
- 純文字 + 形狀，無外部素材相依，可重複執行重新生成。

重新生成：python3 build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# 色彩與字型
# ----------------------------------------------------------------------------
INK      = RGBColor(0x1A, 0x23, 0x32)   # 主文字（深墨藍）
BRAND    = RGBColor(0x0E, 0x2A, 0x47)   # 品牌深藍
ACCENT   = RGBColor(0x00, 0xA8, 0xA8)   # 主強調（青）
MUTED    = RGBColor(0x5B, 0x6B, 0x7B)   # 次要灰
LIGHT    = RGBColor(0xF4, 0xF6, 0xF8)   # 淺底面板
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0xD8, 0xDE, 0xE4)   # 框線
TBD_BG   = RGBColor(0xFF, 0xF3, 0xCD)   # 待填欄位底（淺琥珀）
TBD_TX   = RGBColor(0x9A, 0x73, 0x00)   # 待填文字

# 各產品識別色
C_WOKYIS = RGBColor(0xE8, 0x74, 0x3B)   # 復古橘
C_LUNA   = RGBColor(0x7C, 0x5C, 0xFF)   # 露娜紫
C_SHARGE = RGBColor(0x1F, 0xA2, 0xFF)   # 科技藍

FONT = "Microsoft JhengHei"   # 微軟正黑體；缺字型時由開啟端套用替代

EMU_PER_IN = 914400


# ----------------------------------------------------------------------------
# 低階繪圖 helper
# ----------------------------------------------------------------------------
def _set_run_font(run, size, color, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # 同步設定東亞字型，確保中文正確套用
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
             shape=MSO_SHAPE.RECTANGLE, dashed=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
        if dashed:
            ln = sp.line._get_or_add_ln()
            d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
            ln.append(d)
    return sp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, space_after=2, wrap=True):
    """runs: list of paragraphs, each para = list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (text, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            _set_run_font(r, size, color, bold)
    return tb


# ----------------------------------------------------------------------------
# 版面框架
# ----------------------------------------------------------------------------
SW, SH = 13.333, 7.5   # 16:9 inches


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def content_header(slide, kicker, title, accent=ACCENT):
    """內頁頁首：左側色條 + kicker + 標題。"""
    add_rect(slide, 0, 0, SW, 1.32, fill=WHITE)
    add_rect(slide, 0.55, 0.42, 0.12, 0.62, fill=accent)
    add_text(slide, 0.8, 0.34, 11.8, 0.85, [
        [(kicker, 11, accent, True)],
        [(title, 24, BRAND, True)],
    ], line_spacing=1.0, space_after=2)
    add_rect(slide, 0.55, 1.28, SW - 1.1, 0.012, fill=LINE)


def footer(slide, page_no):
    add_text(slide, 0.55, SH - 0.5, 6.0, 0.3,
             [[("LYCANDER × Wokyis × Sharge ｜ 8/1 三檔齊發團購提報", 8.5, MUTED, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - 1.4, SH - 0.5, 0.85, 0.3,
             [[(str(page_no), 8.5, MUTED, False)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def tag_chip(slide, x, y, text, color):
    w = 0.28 + 0.108 * len(text)
    add_rect(slide, x, y, w, 0.34, fill=None, line=color, line_w=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x, y, w, 0.34, [[(text, 9.5, color, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


# ----------------------------------------------------------------------------
# 各類頁面
# ----------------------------------------------------------------------------
PRODUCTS = [
    {
        "key": "Wokyis M5",
        "brand": "Wokyis",
        "color": C_WOKYIS,
        "tagline": "13 合 1 迷你螢幕擴充座 · 復古外觀，現代武裝",
        "intro": ("專為 Mac mini 打造的 13 合 1 迷你螢幕擴充座，復古外觀內藏現代高速規格；"
                  "同時相容 MacBook Air / Pro 與 Windows 筆電、桌機，桌面一站到位。"),
        "chips": ["13 合 1 擴充", "復古造型", "雙版本 10G/80G", "內建螢幕"],
        "highlights": [
            "13 合 1 多埠擴充，桌面一線整合到位",
            "復古造型 × 現代效能，兼具桌面質感與實用",
            "雙版本：標準版 10Gbps／專業版 80Gbps Thunderbolt 5 Dock",
            "內建螢幕顯示連接狀態，資訊一目了然",
            "完整配件：雙長度 USB-C 線、散熱墊、螺絲起子、快速指南",
        ],
        "specs": [
            ("版本", "標準版 10Gbps／專業版 80Gbps（Thunderbolt 5 Dock）"),
            ("連接埠", "13 合 1（HDMI／USB-A／USB-C／讀卡機 等，依版本）"),
            ("相容", "Mac mini、MacBook Air / Pro、Windows 筆電 / 桌機"),
            ("隨附配件", "18mm & 100mm USB-C 對 USB-C 線、供電線、散熱墊、螺絲起子、快速指南"),
            ("適用情境", "桌面整線、外接高速 SSD、影音剪輯／文書處理"),
        ],
        "url": "https://www.zeczec.com/projects/wokyis",
    },
    {
        "key": "LUNA Mag 小露娜",
        "brand": "LYCANDER",
        "color": C_LUNA,
        "tagline": "口袋型無線充電站 · 一秒展開電就到",
        "intro": ("口袋型 3-in-1 無線充電站，一秒展開即可同時為手機、耳機、手錶充電；"
                  "採次世代 Qi2.2 標準與 MagSafe 磁吸對位，外出隨身的充電解決方案。"),
        "chips": ["3-in-1", "Qi2.2", "MagSafe", "口袋型"],
        "highlights": [
            "3-in-1 同時充手機／耳機／手錶，一站解決",
            "次世代 Qi2.2 標準，充電速度較前代提升約 70%",
            "MagSafe 磁吸精準對位，放上即充",
            "口袋型輕巧折疊，一秒展開即用",
            "差旅／外出隨身充電站，桌面零雜線",
        ],
        "specs": [
            ("充電協定", "Qi2.2，MagSafe 相容"),
            ("充電裝置", "手機 + 耳機 + 手錶（3-in-1 同時）"),
            ("形式", "口袋型可折疊無線充電站，一秒展開"),
            ("核心訴求", "輕便易攜、快速展開、桌面整潔"),
            ("適用情境", "通勤、差旅、辦公桌、床頭隨身充電"),
        ],
        "url": "https://www.zeczec.com/projects/luna-mag",
    },
    {
        "key": "Sharge Disk Pro",
        "brand": "Sharge",
        "color": C_SHARGE,
        "tagline": "全球首款主動散熱 PSSD × 多埠 Hub",
        "intro": ("全球首款內建 NVMe PSSD 的主動散熱多埠 Hub。升級渦輪風扇有效抑制過熱與降速，"
                  "5-in-1 一機整合儲存與擴充；口袋尺寸、重量僅約 150g。"),
        "chips": ["主動散熱 PSSD", "5-in-1 Hub", "最高 4TB", "150g"],
        "highlights": [
            "全球首款主動散熱 PSSD + Hub，Ice-storm 散熱抑制降速",
            "5-in-1 多埠：4 埠 Hub + HDMI 2.1 + 內建線",
            "最高 4TB 容量、10Gbps 高速傳輸",
            "MagSafe 磁吸、約 150g 口袋尺寸",
            "集結 3000+ 前代募資者回饋優化（散熱／HDMI 穩定／供電／相容／埠數）",
        ],
        "specs": [
            ("容量", "最高 4TB"),
            ("傳輸速度", "10Gbps"),
            ("介面", "4 埠 Hub、HDMI 2.1、內建線、MagSafe 磁吸"),
            ("散熱", "Ice-storm 主動散熱渦輪風扇"),
            ("重量", "約 150g（口袋尺寸）"),
        ],
        "url": "https://www.zeczec.com/projects/disk-pro",
    },
]

# 團購商業條件欄位（全部待填，不臆測）
OFFER_FIELDS = ["建議售價 (NT$)", "募資價 (NT$)", "團購價 (NT$)", "拆帳 %", "MOQ 起訂量", "出貨日"]
TBD = "待填 TBD"


def slide_cover(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SW, SH, fill=BRAND)
    add_rect(s, 0, 0, SW, 0.18, fill=ACCENT)
    # 三色識別條
    add_rect(s, 0.9, 2.05, 0.16, 1.0, fill=C_WOKYIS)
    add_rect(s, 1.12, 2.05, 0.16, 1.0, fill=C_LUNA)
    add_rect(s, 1.34, 2.05, 0.16, 1.0, fill=C_SHARGE)

    add_text(s, 0.9, 1.45, 11, 0.5, [[("2026 夏季團購檔期提報", 14, ACCENT, True)]])
    add_text(s, 0.9, 3.15, 11.5, 1.8, [
        [("三檔齊發 ‧ 團購檔期提報", 40, WHITE, True)],
        [("Wokyis × LYCANDER × Sharge", 22, RGBColor(0xC8, 0xD6, 0xE5), True)],
    ], line_spacing=1.05, space_after=8)

    # 資訊列
    add_rect(s, 0.92, 5.35, 5.4, 0.92, fill=RGBColor(0x16, 0x38, 0x5C),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, 1.2, 5.5, 5.0, 0.62, [
        [("上檔日期　", 12, RGBColor(0x9F, 0xB4, 0xC9), False), ("8 / 1 三檔同步上檔", 14, WHITE, True)],
        [("提報對象　", 12, RGBColor(0x9F, 0xB4, 0xC9), False), ("團購業者 ama shop", 13, WHITE, True)],
    ], line_spacing=1.15, space_after=4)

    add_text(s, 0.9, SH - 0.75, 11.5, 0.4,
             [[("LYCANDER 提供　｜　service@lycander.tw", 11, RGBColor(0x9F, 0xB4, 0xC9), False)]])
    return s


def slide_overview(prs, page):
    s = blank_slide(prs)
    content_header(s, "OVERVIEW", "三檔齊發總覽", ACCENT)
    add_text(s, 0.8, 1.5, 11.7, 0.5,
             [[("三檔嘖嘖募資人氣商品，8/1 同步上檔，一次提供團購主完整 3C 選品組合。", 13, MUTED, False)]])

    cw, gap, x0, y0, ch = 3.75, 0.28, 0.8, 2.2, 4.0
    for i, p in enumerate(PRODUCTS):
        x = x0 + i * (cw + gap)
        add_rect(s, x, y0, cw, ch, fill=WHITE, line=LINE, line_w=1.0,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(s, x, y0, cw, 0.16, fill=p["color"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 圖片佔位
        add_rect(s, x + 0.3, y0 + 0.42, cw - 0.6, 1.5, fill=LIGHT, line=LINE,
                 line_w=1.0, dashed=True)
        add_text(s, x + 0.3, y0 + 0.42, cw - 0.6, 1.5,
                 [[("產品圖 / KV 置入處", 10, MUTED, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.3, y0 + 2.05, cw - 0.6, 0.35,
                 [[(p["brand"], 10, p["color"], True)]])
        add_text(s, x + 0.3, y0 + 2.34, cw - 0.6, 0.5,
                 [[(p["key"], 16, BRAND, True)]])
        add_text(s, x + 0.3, y0 + 2.92, cw - 0.6, 0.85,
                 [[(p["tagline"], 11, INK, False)]], line_spacing=1.1)
        add_text(s, x + 0.3, y0 + ch - 0.42, cw - 0.6, 0.32,
                 [[("嘖嘖募資專案 ▸", 9.5, p["color"], True)]])
    footer(s, page)
    return s


def slide_proposal(prs, page):
    s = blank_slide(prs)
    content_header(s, "PROPOSAL", "提案概要：為何三檔齊發", ACCENT)
    blocks = [
        ("一次到位的 3C 選品", "擴充座 × 無線充電 × 行動儲存，覆蓋 Mac／行動辦公族的桌面與外出需求，提高客單與連帶銷售。"),
        ("募資人氣背書", "三檔皆為嘖嘖募資專案，具話題與聲量基礎，導購素材完整、轉化信任度高。"),
        ("目標客群高度重疊", "Apple／Mac 生態、科技嘗鮮族、行動辦公與差旅族群；一次團購觸及同一群高消費力受眾。"),
        ("齊發檔期綜效", "8/1 三檔同步上檔，可組合加購、滿額優惠，拉高團單金額與開團聲量。"),
    ]
    cw, gap, x0, y0, ch = 5.75, 0.35, 0.8, 1.95, 2.05
    for i, (h, d) in enumerate(blocks):
        x = x0 + (i % 2) * (cw + gap)
        y = y0 + (i // 2) * (ch + 0.3)
        add_rect(s, x, y, cw, ch, fill=LIGHT, line=LINE, line_w=1.0,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(s, x, y + 0.28, 0.1, ch - 0.56, fill=ACCENT)
        add_text(s, x + 0.35, y + 0.28, cw - 0.6, 0.45, [[(h, 15, BRAND, True)]])
        add_text(s, x + 0.35, y + 0.86, cw - 0.65, ch - 1.0,
                 [[(d, 12, INK, False)]], line_spacing=1.18)
    # 團購主可得支援
    add_text(s, 0.8, y0 + 2 * ch + 0.45, 11.7, 0.4,
             [[("團購主可獲支援：", 12, ACCENT, True),
               ("產品圖／KV／規格表／賣點文案／QA 懶人包／開團素材包（依檔期提供）", 12, INK, False)]])
    footer(s, page)
    return s


def slide_product_intro(prs, p, page):
    s = blank_slide(prs)
    content_header(s, p["brand"].upper(), p["key"], p["color"])
    # 左：圖片佔位
    add_rect(s, 0.8, 1.7, 4.7, 4.4, fill=LIGHT, line=LINE, line_w=1.0, dashed=True)
    add_text(s, 0.8, 1.7, 4.7, 4.4,
             [[("產品圖 / KV", 13, MUTED, True)], [("置入處", 11, MUTED, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 右：tagline + intro + chips + highlights
    rx = 5.85
    add_text(s, rx, 1.7, 6.7, 0.55, [[(p["tagline"], 15, p["color"], True)]],
             line_spacing=1.05)
    add_text(s, rx, 2.3, 6.7, 0.95, [[(p["intro"], 12, INK, False)]], line_spacing=1.25)
    # chips
    cx = rx
    for c in p["chips"]:
        w = tag_chip(s, cx, 3.35, c, p["color"])
        cx += w + 0.18
    # highlights
    add_text(s, rx, 3.95, 6.7, 0.35, [[("核心賣點", 12, BRAND, True)]])
    para = []
    for h in p["highlights"]:
        para.append([("▸  ", 12, p["color"], True), (h, 12, INK, False)])
    add_text(s, rx, 4.35, 6.8, 1.9, para, line_spacing=1.15, space_after=5)
    footer(s, page)
    return s


def slide_product_offer(prs, p, page):
    s = blank_slide(prs)
    content_header(s, p["brand"].upper(), p["key"] + "　規格與團購方案", p["color"])
    # 左：規格表
    add_text(s, 0.8, 1.65, 5.8, 0.35, [[("產品規格", 13, BRAND, True)]])
    y = 2.1
    for k, v in p["specs"]:
        add_rect(s, 0.8, y, 1.7, 0.62, fill=LIGHT, line=LINE, line_w=0.75)
        add_text(s, 0.92, y, 1.5, 0.62, [[(k, 11, BRAND, True)]],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, 2.5, y, 4.1, 0.62, fill=WHITE, line=LINE, line_w=0.75)
        add_text(s, 2.64, y, 3.85, 0.62, [[(v, 10.5, INK, False)]],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y += 0.62

    # 右：團購方案（待填）
    add_text(s, 6.95, 1.65, 5.6, 0.35, [[("團購方案（待業務填寫）", 13, p["color"], True)]])
    y2 = 2.1
    for f in OFFER_FIELDS:
        add_rect(s, 6.95, y2, 2.6, 0.62, fill=LIGHT, line=LINE, line_w=0.75)
        add_text(s, 7.1, y2, 2.4, 0.62, [[(f, 11, BRAND, True)]],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, 9.55, y2, 3.0, 0.62, fill=TBD_BG, line=LINE, line_w=0.75)
        add_text(s, 9.7, y2, 2.8, 0.62, [[(TBD, 11, TBD_TX, True)]],
                 anchor=MSO_ANCHOR.MIDDLE)
        y2 += 0.62
    add_text(s, 6.95, y2 + 0.05, 5.6, 0.35,
             [[("嘖嘖募資專案：" + p["url"], 9.5, MUTED, False)]])
    footer(s, page)
    return s


def slide_mechanics(prs, page):
    s = blank_slide(prs)
    content_header(s, "CAMPAIGN", "團購機制與檔期時程", ACCENT)
    # 時程條
    add_text(s, 0.8, 1.55, 11.7, 0.35, [[("檔期時程（8/1 三檔同步上檔）", 13, BRAND, True)]])
    steps = [("預熱", "上檔前\n素材／名單預告"), ("開團", "8 / 1\n三檔齊發上線"),
             ("加溫", "檔中\n滿額／加購推播"), ("結團", "待填 TBD\n統計與請款"),
             ("出貨", "待填 TBD\n依各檔到貨")]
    n = len(steps)
    cw, gap, x0, y = 2.18, 0.18, 0.8, 2.05
    for i, (t, d) in enumerate(steps):
        x = x0 + i * (cw + gap)
        add_rect(s, x, y, cw, 1.2, fill=BRAND if i == 1 else LIGHT,
                 line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tcol = WHITE if i == 1 else BRAND
        dcol = RGBColor(0xC8, 0xD6, 0xE5) if i == 1 else MUTED
        add_text(s, x + 0.15, y + 0.16, cw - 0.3, 0.4, [[(t, 14, tcol, True)]],
                 align=PP_ALIGN.CENTER)
        add_text(s, x + 0.12, y + 0.6, cw - 0.24, 0.55,
                 [[(line, 9.5, dcol, False)] for line in d.split("\n")],
                 align=PP_ALIGN.CENTER, line_spacing=1.0)

    # 機制重點
    add_text(s, 0.8, 3.6, 11.7, 0.35, [[("團購機制重點", 13, BRAND, True)]])
    items = [
        "三檔齊發：8/1 同步上檔，集中聲量、一次開團三主打。",
        "組合優惠：擴充座 × 充電 × 儲存可設加購／滿額門檻，拉高客單（門檻待定 TBD）。",
        "素材支援：提供產品圖、KV、規格表、賣點文案、QA 懶人包等開團素材包。",
        "結團與出貨：結團日、請款與各檔出貨時程待雙方確認（TBD）。",
    ]
    para = [[("▸  ", 12, ACCENT, True), (t, 12, INK, False)] for t in items]
    add_text(s, 0.8, 4.0, 11.9, 2.0, para, line_spacing=1.2, space_after=7)
    footer(s, page)
    return s


def slide_summary_table(prs, page):
    s = blank_slide(prs)
    content_header(s, "TERMS", "合作條件總表（待填 TBD）", ACCENT)
    add_text(s, 0.8, 1.5, 11.7, 0.4,
             [[("以下商業條件待雙方確認後填入，本提報先提供完整商品與活動資訊。", 12, MUTED, False)]])

    rows = ["項目"] + OFFER_FIELDS + ["結帳/請款", "物流", "售後保固"]
    cols = ["", "Wokyis M5", "LUNA Mag 小露娜", "Sharge Disk Pro"]
    x0, y0 = 0.8, 2.05
    label_w, col_w, rh = 2.5, 3.0, 0.42

    for ri, rname in enumerate(rows):
        y = y0 + ri * rh
        # 標籤欄
        is_head = ri == 0
        add_rect(s, x0, y, label_w, rh, fill=BRAND if is_head else LIGHT,
                 line=LINE, line_w=0.5)
        add_text(s, x0 + 0.12, y, label_w - 0.2, rh,
                 [[(rname, 10.5, WHITE if is_head else BRAND, True)]],
                 anchor=MSO_ANCHOR.MIDDLE)
        for ci in range(3):
            x = x0 + label_w + ci * col_w
            if is_head:
                add_rect(s, x, y, col_w, rh, fill=PRODUCTS[ci]["color"], line=LINE, line_w=0.5)
                add_text(s, x + 0.1, y, col_w - 0.2, rh,
                         [[(cols[ci + 1], 11, WHITE, True)]],
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                add_rect(s, x, y, col_w, rh, fill=TBD_BG, line=LINE, line_w=0.5)
                add_text(s, x + 0.1, y, col_w - 0.2, rh, [[(TBD, 10, TBD_TX, True)]],
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


def slide_closing(prs, page):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SW, SH, fill=BRAND)
    add_rect(s, 0, 0, SW, 0.18, fill=ACCENT)
    add_text(s, 0.9, 2.15, 11.5, 1.4, [
        [("一次提案，三檔齊發", 34, WHITE, True)],
        [("8/1 同步上檔，誠摯邀請 ama shop 共同開團", 16, RGBColor(0xC8, 0xD6, 0xE5), False)],
    ], line_spacing=1.1, space_after=8)

    add_text(s, 0.92, 3.95, 11, 0.4, [[("聯絡窗口", 12, ACCENT, True)]])
    add_text(s, 0.92, 4.3, 11, 0.45, [
        [("LYCANDER　｜　service@lycander.tw", 15, WHITE, True)]])

    add_text(s, 0.92, 5.15, 11.5, 1.4, [
        [("Wokyis M5　", 12, C_WOKYIS, True), ("https://www.zeczec.com/projects/wokyis", 11, RGBColor(0xC8, 0xD6, 0xE5), False)],
        [("LUNA Mag　 ", 12, C_LUNA, True), ("https://www.zeczec.com/projects/luna-mag", 11, RGBColor(0xC8, 0xD6, 0xE5), False)],
        [("Disk Pro　  ", 12, C_SHARGE, True), ("https://www.zeczec.com/projects/disk-pro", 11, RGBColor(0xC8, 0xD6, 0xE5), False)],
    ], line_spacing=1.25, space_after=4)
    return s


# ----------------------------------------------------------------------------
# 組裝
# ----------------------------------------------------------------------------
def build(path="三檔齊發團購提報_0801.pptx"):
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_PER_IN))
    prs.slide_height = Emu(int(SH * EMU_PER_IN))

    page = 1
    slide_cover(prs)                       # 1 封面（不標頁碼）
    slide_overview(prs, (page := page + 1))    # 2
    slide_proposal(prs, (page := page + 1))    # 3
    for p in PRODUCTS:                          # 4-9
        slide_product_intro(prs, p, (page := page + 1))
        slide_product_offer(prs, p, (page := page + 1))
    slide_mechanics(prs, (page := page + 1))   # 10
    slide_summary_table(prs, (page := page + 1))  # 11
    slide_closing(prs, page + 1)               # 12 結尾

    prs.save(path)
    print(f"已生成：{path}（{len(prs.slides.__iter__.__self__._sldIdLst)} 頁）")
    return path


if __name__ == "__main__":
    import os
    here = os.path.dirname(os.path.abspath(__file__))
    build(os.path.join(here, "三檔齊發團購提報_0801.pptx"))
