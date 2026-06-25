# -*- coding: utf-8 -*-
"""德誼 Roadshow 活動簡報產生器 (Wokyis × LYCANDER) — 可編輯 PPTX, 繁體中文, 16:9"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import sys

# ---------- 品牌色彩 ----------
INK    = RGBColor(0x1D, 0x1D, 0x1F)   # 石墨黑（主文字/深底）
SLATE  = RGBColor(0x6E, 0x6E, 0x73)   # 中灰（次要文字）
MIST   = RGBColor(0x9A, 0x9A, 0xA0)   # 淺灰
CLOUD  = RGBColor(0xF5, 0xF5, 0xF7)   # 雲灰（淺底）
LINE   = RGBColor(0xE3, 0xE3, 0xE8)   # 分隔線
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FLAME  = RGBColor(0xFF, 0x5A, 0x36)   # 焰橙（快閃促銷主色）
FLAME_D= RGBColor(0xD8, 0x3F, 0x1E)
AZURE  = RGBColor(0x0A, 0x84, 0xFF)   # 科技藍（LYCANDER）
GOLD   = RGBColor(0xB8, 0x89, 0x3B)   # 金（premium 贈品）
INKSOFT= RGBColor(0x2C, 0x2C, 0x2E)

CJK = "Microsoft JhengHei"            # 微軟正黑體（繁中安全字）
CJK_L = "Microsoft JhengHei Light"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
EMUW = prs.slide_width
EMUH = prs.slide_height

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMUW, EMUH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def _set_cjk(run, font):
    """讓中文也套用指定字型（east asian）"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, space_after=0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,font,charspace)
       簡寫: 也接受單一 list-of-tuples 視為一段"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for seg in para:
            txt, size, color, bold, font = seg[0], seg[1], seg[2], seg[3], seg[4]
            cs = seg[5] if len(seg) > 5 else None
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
            _set_cjk(r, font)
            if cs is not None:
                r.font._rPr.set('spc', str(int(cs)))
    return tb

def chip(s, x, y, w, h, label, fill, fg=WHITE, size=12, bold=True, radius=True):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    c = rect(s, x, y, w, h, fill=fill, shape=shp)
    try:
        c.adjustments[0] = 0.5
    except Exception:
        pass
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left=Inches(0.06); tf.margin_right=Inches(0.06)
    tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=fg; r.font.name=CJK
    _set_cjk(r, CJK)
    return c

MX = 0.92   # 左右安全邊界

def header(s, kicker, title, accent=FLAME, page=None, total=12):
    rect(s, MX, 0.62, 0.07, 0.42, fill=accent)
    text(s, MX+0.2, 0.55, 9.5, 0.4, [(kicker, 12.5, accent, True, CJK, 220)])
    text(s, MX+0.2, 0.9, 11.0, 0.85, [(title, 30, INK, True, CJK)])
    rect(s, MX, 1.86, 13.333-2*MX, 0.018, fill=LINE)
    if page:
        text(s, 13.333-MX-1.6, 0.62, 1.6, 0.4,
             [[(f"{page:02d}", 12, INK, True, CJK), (f" / {total}", 12, MIST, False, CJK)]],
             align=PP_ALIGN.RIGHT)

def footer(s, dark=False):
    c = MIST if not dark else RGBColor(0x8A,0x8A,0x90)
    text(s, MX, 7.04, 6.0, 0.3,
         [[("德誼 Roadshow　", 9, c, True, CJK),
           ("Wokyis × LYCANDER 快閃活動簡報", 9, c, False, CJK)]])
    text(s, 13.333-MX-4.0, 7.04, 4.0, 0.3,
         [[("機密・限行銷協作使用", 9, c, False, CJK)]], align=PP_ALIGN.RIGHT)

# =================================================================
# 01 — 封面
# =================================================================
s = slide(INK)
rect(s, 0, 0, 13.333, 7.5, fill=INK)
# 右側裝飾色塊
rect(s, 9.6, 0, 3.733, 7.5, fill=INKSOFT)
rect(s, 9.6, 0, 0.07, 7.5, fill=FLAME)
# 焰橙幾何
rect(s, 10.4, 1.0, 2.1, 2.1, fill=FLAME, shape=MSO_SHAPE.OVAL)
rect(s, 10.9, 3.7, 1.5, 1.5, fill=None, line=AZURE, line_w=2.2, shape=MSO_SHAPE.OVAL)
chip(s, MX, 1.15, 3.2, 0.5, "2026 夏季快閃巡迴 ROADSHOW", FLAME, size=12)
text(s, MX, 1.95, 8.6, 2.4,
     [[("德誼 Roadshow", 60, WHITE, True, CJK)],
      [("活動企劃簡報", 60, WHITE, True, CJK)]],
     line_spacing=1.02)
text(s, MX, 4.35, 8.4, 0.6,
     [[("主題品牌　", 16, MIST, False, CJK),
       ("Wokyis", 16, WHITE, True, CJK),
       ("　×　", 16, MIST, False, CJK),
       ("LYCANDER", 16, WHITE, True, CJK)]])
# 場次卡
rect(s, MX, 5.15, 3.7, 1.15, fill=INKSOFT)
rect(s, MX, 5.15, 0.06, 1.15, fill=FLAME)
text(s, MX+0.28, 5.3, 3.3, 0.4, [[("場次一　", 12, MIST, False, CJK),("7 / 5", 22, WHITE, True, CJK)]])
text(s, MX+0.28, 5.82, 3.3, 0.4, [[("新竹巨城 Big City", 14, WHITE, False, CJK)]])
rect(s, MX+4.0, 5.15, 3.7, 1.15, fill=INKSOFT)
rect(s, MX+4.0, 5.15, 0.06, 1.15, fill=AZURE)
text(s, MX+4.28, 5.3, 3.3, 0.4, [[("場次二　", 12, MIST, False, CJK),("7 / 12", 22, WHITE, True, CJK)]])
text(s, MX+4.28, 5.82, 3.3, 0.4, [[("台北三創 Syntrend", 14, WHITE, False, CJK)]])
text(s, MX, 6.75, 8.0, 0.4,
     [[("交付對象：行銷協作公司　|　提案版本 v1　|　日期：2026.06.25", 11, MIST, False, CJK)]])

# =================================================================
# 02 — 活動總覽 (一頁看懂)
# =================================================================
s = slide(WHITE)
header(s, "ACTIVITY OVERVIEW", "活動總覽・一頁看懂", FLAME, page=2)
cards = [
    ("活動性質", "通路快閃 Roadshow", "於德誼門市／百貨設櫃，主打體驗＋限時優惠", FLAME),
    ("檔期場次", "2 場・7 月", "7/5 新竹巨城　|　7/12 台北三創", AZURE),
    ("主題品牌", "Wokyis × LYCANDER", "Wokyis M5 為主推英雄商品，LYCANDER 周邊加值", INK),
    ("核心優惠", "快閃 9 折", "活動期間全品項限時 9 折，營造搶購氛圍", FLAME),
    ("加碼贈品", "買 M5 雙重送", "Mac Mini 鋁合金支架＋HDMI 2.1 8K 影音線", GOLD),
    ("滿額好禮", "單筆滿 $1,000", "贈 YOYOISLES AirTag 保護套（市值 $190）", AZURE),
]
cw, ch, gx, gy = 3.72, 1.92, 0.21, 0.22
x0, y0 = MX, 2.18
for i, (tag, big, desc, ac) in enumerate(cards):
    cx = x0 + (i % 3) * (cw + gx)
    cy = y0 + (i // 3) * (ch + gy)
    rect(s, cx, cy, cw, ch, fill=CLOUD)
    rect(s, cx, cy, 0.06, ch, fill=ac)
    text(s, cx+0.26, cy+0.2, cw-0.5, 0.3, [(tag, 11.5, ac, True, CJK, 120)])
    text(s, cx+0.26, cy+0.5, cw-0.45, 0.55, [(big, 19, INK, True, CJK)])
    text(s, cx+0.26, cy+1.08, cw-0.45, 0.7, [(desc, 11.5, SLATE, False, CJK)], line_spacing=1.12)
footer(s)

# =================================================================
# 03 — 活動目的與目標
# =================================================================
s = slide(WHITE)
header(s, "OBJECTIVES", "活動目的與策略目標", FLAME, page=3)
goals = [
    ("01", "推升 Wokyis M5 銷量", "以快閃 9 折＋雙重贈品打造「現在買最划算」的急迫感，集中火力衝刺英雄商品。"),
    ("02", "建立品牌體驗", "於高人流地點讓消費者實際接觸 Wokyis 與 LYCANDER，強化品牌記憶與信任。"),
    ("03", "帶動連帶銷售", "用「單筆滿千送」拉高客單價，將 LYCANDER 周邊與配件一併帶出。"),
    ("04", "蒐集名單與聲量", "現場互動、社群打卡與會員加入，累積後續再行銷的第一方名單。"),
]
gw = (13.333-2*MX-0.3)/2
gh = 1.95
for i,(n,t,d) in enumerate(goals):
    gx = MX + (i%2)*(gw+0.3)
    gy = 2.2 + (i//2)*(gh+0.25)
    rect(s, gx, gy, gw, gh, fill=CLOUD)
    text(s, gx+0.32, gy+0.26, 1.2, 0.9, [(n, 40, FLAME, True, CJK)])
    text(s, gx+1.55, gy+0.32, gw-1.85, 0.5, [(t, 18, INK, True, CJK)])
    text(s, gx+1.55, gy+0.86, gw-1.9, 0.9, [(d, 12.5, SLATE, False, CJK)], line_spacing=1.2)
footer(s)

# =================================================================
# 04 — 檔期與場地
# =================================================================
s = slide(WHITE)
header(s, "SCHEDULE & VENUE", "活動檔期與場地", AZURE, page=4)
def venue(x, ac, day, place, sub, rows):
    w = (13.333-2*MX-0.34)/2
    rect(s, x, 2.18, w, 4.5, fill=CLOUD)
    rect(s, x, 2.18, w, 0.9, fill=ac)
    text(s, x+0.34, 2.32, w-1.0, 0.5, [[(day, 30, WHITE, True, CJK),("（日）", 14, WHITE, False, CJK)]])
    text(s, x+0.34, 2.86, w-0.6, 0.3, [[(place, 14, WHITE, True, CJK),("　",12,WHITE,False,CJK),(sub,12,WHITE,False,CJK)]])
    yy = 3.35
    for label, val in rows:
        text(s, x+0.34, yy, 2.0, 0.4, [(label, 12.5, SLATE, True, CJK)])
        text(s, x+2.2, yy, w-2.6, 0.55, [(val, 13.5, INK, False, CJK)], line_spacing=1.1)
        yy += 0.66
        rect(s, x+0.34, yy-0.12, w-0.68, 0.012, fill=LINE)
venue(MX, FLAME, "7 / 5", "新竹巨城購物中心", "Big City", [
    ("地點", "新竹巨城・人流主動線設櫃"),
    ("建議時段", "11:00 – 21:00（百貨營業時間）"),
    ("主打", "Wokyis M5 體驗 × 快閃 9 折"),
    ("客群", "新竹科學園區家庭、3C 重度族群"),
    ("佈場", "活動前一日 18:00 後進場佈置"),
])
venue(MX+(13.333-2*MX-0.34)/2+0.34, AZURE, "7 / 12", "台北三創生活園區", "Syntrend", [
    ("地點", "三創・3C 指標賣場樓層"),
    ("建議時段", "11:00 – 21:30"),
    ("主打", "Wokyis M5 × LYCANDER 周邊組合"),
    ("客群", "北市科技嘗鮮族、Apple 生態用戶"),
    ("佈場", "活動前一日 19:00 後進場佈置"),
])
footer(s)

# =================================================================
# 05 — 主題品牌
# =================================================================
s = slide(WHITE)
header(s, "FEATURED BRANDS", "主題品牌與商品線", INK, page=5)
# Wokyis
rect(s, MX, 2.18, 5.86, 4.5, fill=INK)
rect(s, MX, 2.18, 5.86, 0.07, fill=FLAME)
text(s, MX+0.4, 2.55, 5.0, 0.6, [("Wokyis", 30, WHITE, True, CJK)])
chip(s, MX+0.4, 3.2, 1.7, 0.42, "英雄主推商品", FLAME, size=11)
text(s, MX+0.4, 3.85, 5.1, 0.5, [[("Wokyis ", 16, WHITE, True, CJK),("M5", 26, FLAME, True, CJK)]])
text(s, MX+0.4, 4.5, 5.1, 1.9, [
    [("・本檔活動的主角，搭配快閃 9 折最有感", 12.5, CLOUD, False, CJK)],
    [("・購買即啟動「雙重贈品」加碼機制", 12.5, CLOUD, False, CJK)],
    [("・現場設體驗區，主打實機操作與情境展示", 12.5, CLOUD, False, CJK)],
    [("・銷售與名單蒐集主要承載品項", 12.5, CLOUD, False, CJK)],
], line_spacing=1.5)
# LYCANDER
rx = MX+6.2
rect(s, rx, 2.18, 5.86, 4.5, fill=CLOUD)
rect(s, rx, 2.18, 5.86, 0.07, fill=AZURE)
text(s, rx+0.4, 2.55, 5.0, 0.6, [("LYCANDER", 30, INK, True, CJK)])
chip(s, rx+0.4, 3.2, 1.9, 0.42, "Apple 周邊配件", AZURE, size=11)
text(s, rx+0.4, 3.85, 5.1, 0.5, [("精品級擴充與影音配件", 16, INK, True, CJK)])
text(s, rx+0.4, 4.5, 5.1, 1.9, [
    [("・Mac Mini M4 & M4 Pro 鋁合金支架", 12.5, INK, True, CJK)],
    [("　質感擴充、與機身一體成型設計", 11.5, SLATE, False, CJK)],
    [("・HDMI 2.1 8K 高畫質認證影音傳輸線（1.5M）", 12.5, INK, True, CJK)],
    [("　支援 8K 極細認證，影音玩家首選", 11.5, SLATE, False, CJK)],
], line_spacing=1.42)
footer(s)

# =================================================================
# 06 — 核心優惠：快閃 9 折
# =================================================================
s = slide(INK)
rect(s, 0,0,13.333,7.5, fill=INK)
rect(s, MX, 0.62, 0.07, 0.42, fill=FLAME)
text(s, MX+0.2, 0.55, 9, 0.4, [("CORE OFFER", 12.5, FLAME, True, CJK, 220)])
text(s, MX+0.2, 0.9, 11, 0.85, [("核心優惠　快閃 9 折", 30, WHITE, True, CJK)])
rect(s, MX, 1.86, 13.333-2*MX, 0.018, fill=INKSOFT)
# 巨大 9折
rect(s, MX, 2.5, 5.3, 3.7, fill=INKSOFT)
rect(s, MX, 2.5, 0.08, 3.7, fill=FLAME)
text(s, MX, 2.75, 5.3, 2.4, [[("9", 150, FLAME, True, CJK),("折", 54, WHITE, True, CJK)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, MX, 5.55, 5.3, 0.5, [("活動期間・限時快閃價", 16, CLOUD, False, CJK)], align=PP_ALIGN.CENTER)
# 說明
bx = MX+5.8
pts = [
    ("適用範圍", "活動現場全品項適用快閃 9 折"),
    ("適用時間", "僅限 7/5、7/12 兩場活動當日"),
    ("搭配機制", "可與「買 M5 雙重贈」「滿千送」並行"),
    ("溝通主軸", "「只有這兩天」── 強化限時急迫感"),
]
yy = 2.5
for t,d in pts:
    rect(s, bx, yy, 6.0, 0.78, fill=INKSOFT)
    rect(s, bx, yy, 0.06, 0.78, fill=FLAME)
    text(s, bx+0.3, yy+0.13, 2.0, 0.5, [(t, 13, FLAME, True, CJK)])
    text(s, bx+2.0, yy+0.13, 4.0, 0.55, [(d, 13, WHITE, False, CJK)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    yy += 0.92
text(s, MX, 6.95, 11, 0.3, [[("※ 折扣與贈品之最終適用條件、是否含特定品項，以德誼公告為準。", 10, MIST, False, CJK)]])

# =================================================================
# 07 — 加碼贈品：買 M5 雙重送
# =================================================================
s = slide(WHITE)
header(s, "BUY M5 · DOUBLE GIFT", "加碼贈品　買 Wokyis M5 雙重送", GOLD, page=7)
chip(s, MX, 2.05, 3.5, 0.5, "購買 Wokyis M5 即享雙重好禮", GOLD, size=12.5)
gifts = [
    ("禮 1", "Mac Mini M4 & M4 Pro Stand", "鋁合金支架", "【LYCANDER】航太級鋁合金，為 Mac Mini M4 / M4 Pro 量身打造的質感擴充底座。", AZURE),
    ("禮 2", "HDMI 2.1 8K 影音傳輸線（1.5M）", "極細認證線材", "【LYCANDER】通過 8K 高畫質認證，極細好收納，支援高更新率影音輸出。", FLAME),
]
gw = (13.333-2*MX-0.3)/2
for i,(badge,name,sub,desc,ac) in enumerate(gifts):
    gx = MX + i*(gw+0.3)
    gy = 2.75
    rect(s, gx, gy, gw, 3.45, fill=CLOUD)
    rect(s, gx, gy, gw, 0.09, fill=ac)
    chip(s, gx+0.34, gy+0.34, 1.05, 0.45, badge, ac, size=13)
    # 圖示占位
    rect(s, gx+gw-1.5, gy+0.3, 1.1, 1.1, fill=WHITE, line=LINE, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, gx+gw-1.5, gy+0.3, 1.1, 1.1, [("商品圖", 11, MIST, False, CJK)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, gx+0.34, gy+1.05, gw-0.7, 0.7, [(name, 17, INK, True, CJK)], line_spacing=1.05)
    text(s, gx+0.34, gy+1.78, gw-0.7, 0.4, [(sub, 13, ac, True, CJK)])
    text(s, gx+0.34, gy+2.25, gw-0.7, 1.0, [(desc, 12, SLATE, False, CJK)], line_spacing=1.25)
text(s, MX, 6.35, 11, 0.4, [[("溝通重點：", 11.5, INK, True, CJK),
    ("一次購買 M5，兩件 LYCANDER 精品配件同步到手，等於「主機＋擴充＋影音」一站到位。", 11.5, SLATE, False, CJK)]])
footer(s)

# =================================================================
# 08 — 滿額贈
# =================================================================
s = slide(WHITE)
header(s, "SPEND & GET", "滿額好禮　單筆滿千即送", AZURE, page=8)
# 左：機制
rect(s, MX, 2.3, 6.4, 4.1, fill=CLOUD)
rect(s, MX, 2.3, 0.07, 4.1, fill=AZURE)
text(s, MX+0.4, 2.65, 5.6, 0.5, [("贈品機制", 13, AZURE, True, CJK, 120)])
text(s, MX+0.4, 3.15, 5.7, 1.0, [[("單筆結帳滿 ", 22, INK, True, CJK),("$1,000", 34, AZURE, True, CJK)]])
text(s, MX+0.4, 4.2, 5.7, 0.5, [("即贈 YOYOISLES AirTag 保護套 1 個", 15, INK, True, CJK)], line_spacing=1.1)
text(s, MX+0.4, 4.95, 5.7, 1.3, [
    [("・以低門檻拉高客單，鼓勵連帶購買", 12.5, SLATE, False, CJK)],
    [("・贈完為止，營造「手刀搶贈」氛圍", 12.5, SLATE, False, CJK)],
    [("・建議現場明示剩餘數量，提升轉換", 12.5, SLATE, False, CJK)],
], line_spacing=1.4)
# 右：贈品卡
rx = MX+6.8
rect(s, rx, 2.3, 4.82, 4.1, fill=INK)
rect(s, rx, 2.3, 4.82, 0.09, fill=GOLD)
rect(s, rx+1.66, 2.7, 1.5, 1.5, fill=INKSOFT, line=GOLD, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+1.66, 2.7, 1.5, 1.5, [("商品圖", 11, MIST, False, CJK)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+0.4, 4.35, 4.0, 0.4, [("【YOYOISLES】", 12, GOLD, True, CJK)], align=PP_ALIGN.CENTER)
text(s, rx+0.4, 4.7, 4.0, 0.7, [("Air Stars AirTag 保護套／保護殼", 15, WHITE, True, CJK)], align=PP_ALIGN.CENTER, line_spacing=1.05)
text(s, rx+0.4, 5.5, 4.0, 0.5, [("門卡磁扣・安全無毒・輕量設計", 11.5, CLOUD, False, CJK)], align=PP_ALIGN.CENTER)
chip(s, rx+1.66, 5.95, 1.5, 0.42, "市值 $190", GOLD, fg=INK, size=12)
footer(s)

# =================================================================
# 09 — 優惠機制總表
# =================================================================
s = slide(WHITE)
header(s, "OFFER SUMMARY", "優惠機制總表・一頁掌握", FLAME, page=9)
cols = ["機制", "條件", "內容 / 贈品", "市值", "備註"]
widths = [2.4, 2.9, 4.2, 1.2, 0.79]
rows = [
    ("快閃 9 折", "活動現場消費", "全品項限時 9 折", "—", "兩場活動日"),
    ("買 M5 贈禮 1", "購買 Wokyis M5", "LYCANDER Mac Mini 鋁合金支架", "—", "雙重送 ①"),
    ("買 M5 贈禮 2", "購買 Wokyis M5", "LYCANDER HDMI 2.1 8K 線（1.5M）", "—", "雙重送 ②"),
    ("滿額贈", "單筆滿 $1,000", "YOYOISLES AirTag 保護套", "$190", "贈完為止"),
]
tx, ty = MX, 2.2
tw = sum(widths)
rowh = 0.92
# 表頭
rect(s, tx, ty, tw, 0.62, fill=INK)
cx = tx
for c,w in zip(cols, widths):
    text(s, cx+0.18, ty, w-0.2, 0.62, [(c, 12.5, WHITE, True, CJK)], anchor=MSO_ANCHOR.MIDDLE)
    cx += w
# 資料列
for i,row in enumerate(rows):
    ry = ty+0.62 + i*rowh
    rect(s, tx, ry, tw, rowh, fill=CLOUD if i%2==0 else WHITE)
    rect(s, tx, ry+rowh, tw, 0.012, fill=LINE)
    rect(s, tx, ry, 0.06, rowh, fill=[FLAME,AZURE,AZURE,GOLD][i])
    cx = tx
    for j,(val,w) in enumerate(zip(row, widths)):
        bold = (j==0)
        col = INK if j!=3 else (GOLD if val!="—" else MIST)
        text(s, cx+0.18, ry, w-0.24, rowh,
             [(val, 12.5 if j!=0 else 13.5, col, bold, CJK)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        cx += w
text(s, MX, ty+0.62+len(rows)*rowh+0.18, 11, 0.4,
     [[("※ 各機制可疊加使用；最終適用條件、品項範圍與庫存以德誼及品牌方公告為準。", 10.5, SLATE, False, CJK)]])
footer(s)

# =================================================================
# 10 — 現場執行重點
# =================================================================
s = slide(WHITE)
header(s, "ON-SITE EXECUTION", "現場執行重點（建議）", AZURE, page=10)
blocks = [
    ("攤位與動線", FLAME, [
        "主視覺背板＋品牌 LOGO 牆",
        "Wokyis M5 體驗區（實機操作）",
        "贈品展示區，清楚標示機制",
        "結帳／名單登錄動線分流",
    ]),
    ("人力配置", AZURE, [
        "現場主管 1 名（總控）",
        "銷售／導購 2–3 名",
        "體驗解說 1 名",
        "名單與打卡引導 1 名",
    ]),
    ("物料清單", GOLD, [
        "主視覺 KV、易拉展、桌卡",
        "優惠／贈品說明牌（含剩餘量）",
        "贈品庫存與簽收表",
        "POS／刷卡與發票設備",
    ]),
    ("互動與蒐集", INK, [
        "社群打卡送小禮 / 抽獎",
        "會員加入 / 名單登錄",
        "現場短影音與照片素材",
        "問卷或意見蒐集",
    ]),
]
bw = (13.333-2*MX-0.66)/4
for i,(t,ac,items) in enumerate(blocks):
    bx = MX + i*(bw+0.22)
    rect(s, bx, 2.18, bw, 4.4, fill=CLOUD)
    rect(s, bx, 2.18, bw, 0.62, fill=ac)
    text(s, bx+0.22, 2.18, bw-0.3, 0.62, [(t, 13.5, WHITE, True, CJK)], anchor=MSO_ANCHOR.MIDDLE)
    yy = 3.0
    for it in items:
        rect(s, bx+0.24, yy+0.12, 0.1, 0.1, fill=ac, shape=MSO_SHAPE.OVAL)
        text(s, bx+0.46, yy, bw-0.6, 0.85, [(it, 11.5, INK, False, CJK)], line_spacing=1.1)
        yy += 0.86
footer(s)

# =================================================================
# 11 — 行銷公司分工與時程
# =================================================================
s = slide(WHITE)
header(s, "ROLES & TIMELINE", "行銷協作分工與時程（建議）", FLAME, page=11)
# 分工表
text(s, MX, 2.1, 6, 0.4, [("協作分工", 14, INK, True, CJK)])
raci = [
    ("品牌方 / LYCANDER", "商品、優惠政策、贈品供應、最終定案", FLAME),
    ("行銷協作公司", "活動規劃、主視覺設計、現場執行、社群曝光", AZURE),
    ("德誼通路", "場地、人流、結帳系統、現場協調", GOLD),
]
yy = 2.6
for who, what, ac in raci:
    rect(s, MX, yy, 5.7, 1.05, fill=CLOUD)
    rect(s, MX, yy, 0.06, 1.05, fill=ac)
    text(s, MX+0.3, yy+0.16, 5.2, 0.4, [(who, 13.5, INK, True, CJK)])
    text(s, MX+0.3, yy+0.56, 5.2, 0.45, [(what, 11.5, SLATE, False, CJK)], line_spacing=1.05)
    yy += 1.2
# 時程
text(s, MX+6.4, 2.1, 6, 0.4, [("關鍵時程", 14, INK, True, CJK)])
timeline = [
    ("即日起", "確認優惠／贈品數量、簽訂協作範圍"),
    ("6 月底前", "主視覺定稿、物料發包、社群預熱規劃"),
    ("活動前 1 週", "物料到位、人力排班、教育訓練"),
    ("7/5・7/12", "現場執行、即時數據回報"),
    ("活動後 1 週", "成效結算、名單交付、檢討報告"),
]
lx = MX+6.6
ly = 2.6
rect(s, lx, ly, 0.025, 3.55, fill=LINE)
for i,(when, what) in enumerate(timeline):
    ny = ly + i*0.72
    rect(s, lx-0.09, ny, 0.22, 0.22, fill=FLAME if i<4 else AZURE, shape=MSO_SHAPE.OVAL)
    text(s, lx+0.35, ny-0.06, 1.7, 0.4, [(when, 12.5, INK, True, CJK)])
    text(s, lx+2.15, ny-0.06, 3.3, 0.6, [(what, 11.5, SLATE, False, CJK)], line_spacing=1.05)
footer(s)

# =================================================================
# 12 — KPI 與聯絡窗口 / 結尾
# =================================================================
s = slide(INK)
rect(s,0,0,13.333,7.5, fill=INK)
rect(s, 9.6, 0, 3.733, 7.5, fill=INKSOFT)
rect(s, 9.6, 0, 0.07, 7.5, fill=FLAME)
rect(s, MX, 0.85, 0.07, 0.42, fill=FLAME)
text(s, MX+0.2, 0.78, 9, 0.4, [("KPI & NEXT STEPS", 12.5, FLAME, True, CJK, 200)])
text(s, MX+0.2, 1.15, 8, 0.7, [("成效指標與後續", 30, WHITE, True, CJK)])
# KPI 卡
kpis = [
    ("英雄商品", "M5 銷售目標", "▲ 衝刺"),
    ("活動客單", "提升連帶率", "↑ 滿千贈"),
    ("名單蒐集", "第一方名單", "＋ 會員"),
    ("品牌聲量", "社群曝光", "＃ 打卡"),
]
for i,(t,big,tag) in enumerate(kpis):
    kx = MX + (i%2)*2.9
    ky = 2.4 + (i//2)*1.5
    rect(s, kx, ky, 2.7, 1.3, fill=INKSOFT)
    rect(s, kx, ky, 0.06, 1.3, fill=FLAME if i%2==0 else AZURE)
    text(s, kx+0.26, ky+0.18, 2.3, 0.3, [(t, 11, MIST, True, CJK)])
    text(s, kx+0.26, ky+0.5, 2.3, 0.4, [(big, 16, WHITE, True, CJK)])
    text(s, kx+0.26, ky+0.92, 2.3, 0.3, [(tag, 12, FLAME if i%2==0 else AZURE, True, CJK)])
text(s, MX, 5.6, 6, 0.4, [[("※ 具體數值目標待三方會議確認後填入。", 10.5, MIST, False, CJK)]])
# 右欄：窗口與下一步
text(s, 9.95, 1.3, 3.0, 0.4, [("下一步", 15, WHITE, True, CJK)])
steps = ["確認優惠與贈品數量","主視覺與物料定稿","排定現場人力與訓練","啟動社群預熱"]
yy=1.9
for stp in steps:
    rect(s, 9.95, yy+0.1, 0.12, 0.12, fill=FLAME, shape=MSO_SHAPE.OVAL)
    text(s, 10.2, yy, 2.9, 0.4, [(stp, 12, CLOUD, False, CJK)])
    yy+=0.52
rect(s, 9.95, 4.2, 2.95, 0.014, fill=RGBColor(0x44,0x44,0x48))
text(s, 9.95, 4.45, 3.0, 0.4, [("聯絡窗口", 13, WHITE, True, CJK)])
text(s, 9.95, 4.9, 3.1, 1.2, [
    [("LYCANDER GROUP", 12, CLOUD, True, CJK)],
    [("service@lycander.tw", 11.5, MIST, False, CJK)],
    [("（行銷協作專案窗口）", 11, MIST, False, CJK)],
], line_spacing=1.35)
text(s, MX, 6.7, 8, 0.4, [[("德誼 Roadshow　Wokyis × LYCANDER　|　2026 夏季快閃巡迴", 11, MIST, False, CJK)]])

out = sys.argv[1] if len(sys.argv) > 1 else "deck.pptx"
prs.save(out)
print("SAVED", out, "slides:", len(prs.slides._sldIdLst))
